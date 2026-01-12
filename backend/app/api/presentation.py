from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Any
from pptx import Presentation
from pptx.util import Inches, Pt
# from pptx.dml.color import RgbColor
# from pptx.enum.text import PP_ALIGN
import io
import base64

router = APIRouter(prefix="/api/presentation", tags=["presentation"])


class SlideContent(BaseModel):
    id: str
    order: int
    title: str
    contentType: str
    content: Any
    notes: str | None = None
    chartConfig: dict[str, Any] | None = None
    chartImage: str | None = None  # Base64 encoded chart image


class PresentationRequest(BaseModel):
    title: str
    slides: list[SlideContent]


@router.post("/generate")
async def generate_presentation(request: PresentationRequest):
    """Generate a PowerPoint file from the presentation configuration."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = request.title

        for slide_data in request.slides:
            if slide_data.contentType == "bullets":
                slide_layout = prs.slide_layouts[1]
            elif slide_data.contentType == "chart":
                slide_layout = prs.slide_layouts[5]
            else:
                slide_layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(slide_layout)

            if slide.shapes.title:
                slide.shapes.title.text = slide_data.title

            if slide_data.contentType == "bullets" and isinstance(slide_data.content, list):
                body_shape = None
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        body_shape = shape
                        break

                if body_shape:
                    tf = body_shape.text_frame
                    tf.clear()
                    for i, bullet in enumerate(slide_data.content):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

            elif slide_data.contentType == "text" and isinstance(slide_data.content, str):
                body_shape = None
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        body_shape = shape
                        break

                if body_shape:
                    body_shape.text_frame.text = slide_data.content

            elif slide_data.contentType == "chart" and slide_data.chartImage:
                image_data = base64.b64decode(slide_data.chartImage)
                image_stream = io.BytesIO(image_data)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(1),
                    Inches(1.5),
                    width=Inches(11),
                    height=Inches(5.5),
                )

            if slide_data.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data.notes

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        filename = f"{request.title.replace(' ', '_')}.pptx"

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


class UpdateSlideRequest(BaseModel):
    slideId: str
    title: str | None = None
    content: Any | None = None
    contentType: str | None = None
    notes: str | None = None


@router.post("/preview")
async def preview_presentation(request: PresentationRequest):
    """Return a preview of the presentation structure without generating the file."""
    preview = {
        "title": request.title,
        "slideCount": len(request.slides),
        "slides": [
            {
                "id": slide.id,
                "order": slide.order,
                "title": slide.title,
                "contentType": slide.contentType,
                "hasChart": slide.chartConfig is not None or slide.chartImage is not None,
            }
            for slide in request.slides
        ],
    }
    return preview
